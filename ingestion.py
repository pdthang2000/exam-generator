import os

import PyPDF2
import chromadb
from openai import OpenAI
from pptx import Presentation

CHROMA_DIR = os.path.join(os.path.dirname(__file__), "chroma_data")


def get_chroma_collection():
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    return client.get_or_create_collection("documents")


def get_openai_client():
    return OpenAI(
        api_key=os.environ.get("OPENAI_API_KEY"),
        base_url=os.environ.get("OPENAI_BASE_URL"),
    )


# --- File parsers ---

def parse_pdf(file_path):
    reader = PyPDF2.PdfReader(file_path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def parse_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


PARSERS = {
    "pdf": parse_pdf,
    "pptx": parse_pptx,
    "txt": parse_txt,
}


def parse_file(file_path, file_type):
    parser = PARSERS.get(file_type)
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(file_path)


# --- Chunking ---

def chunk_text(text, chunk_size=2000, overlap=200):
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))

        if end < len(text):
            break_at = text.rfind("\n\n", start + chunk_size // 2, end)
            if break_at == -1:
                break_at = text.rfind(". ", start + chunk_size // 2, end)
                if break_at != -1:
                    break_at += 2
            if break_at > start:
                end = break_at

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        if end >= len(text):
            break
        start = end - overlap

    return chunks


# --- Ingest / delete ---

def ingest_document(doc_id, text):
    chunks = chunk_text(text)
    if not chunks:
        return 0

    collection = get_chroma_collection()

    ids = [f"doc{doc_id}_chunk{i}" for i in range(len(chunks))]
    metadatas = [{"document_id": doc_id, "chunk_index": i} for i in range(len(chunks))]

    # ChromaDB embeds automatically using its built-in model
    collection.add(
        ids=ids,
        documents=chunks,
        metadatas=metadatas,
    )
    return len(chunks)


def delete_document_chunks(doc_id):
    collection = get_chroma_collection()
    results = collection.get(where={"document_id": doc_id})
    if results["ids"]:
        collection.delete(ids=results["ids"])
